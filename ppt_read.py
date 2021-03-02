import math
import shutil

import win32com.client
import xlrd
import xlwt
from aip import AipSpeech
from moviepy.editor import *
from pptx import Presentation

""" 你的 APPID AK SK """
APP_ID = '22663952'
API_KEY = '5bMW9nKQ6Oxqlvif31mr07h2'
SECRET_KEY = 'GzabHXI5mgFRwhKT2rhsTVNb3xBDtaPb'
#
client = AipSpeech(APP_ID, API_KEY, SECRET_KEY)


def reading(text, name):
	'''
	朗读程序
	:param text: 朗读内容
	:param name: 生成文件名
	:return:
	'''
	print(text, name)
	
	result = client.synthesis(text, 'zh', 1, {
		'spd': 5,
		'vol': 5,
		'pit': 5,
		'per': 0,
	})
	print(text)
	
	if not isinstance(result, dict):
		with open('./audios/%s.mp3' % (str(name)), 'wb') as f:
			f.write(result)


def ppt2jpg():
	'''
	PPT转为JPG图片
	:return:
	'''
	ppt_app = win32com.client.Dispatch('PowerPoint.Application')
	ppt_app.Visible = 1
	print(os.getcwd() + '\朗读.pptx')
	ppt = ppt_app.Presentations.Open(os.getcwd() + '\朗读.pptx')  # 打开 ppt
	
	ppt.SaveAs(os.getcwd() + "\images", 17)  # 17数字是转为 ppt 转为图片
	# # ppt.SaveAs("D:\PythonTest\ppt_read\images", 17)  # 17数字是转为 ppt 转为图片
	print(len(ppt.Slides))
	ppt_app.Quit()  # 关闭资源，退出


def ppt2xls():
	'''
	PPT文字转为EXCEL文件
	:return:
	'''
	prs = Presentation('朗读.pptx')
	# 获取slide幻灯片
	picture_num = 0
	count_num = 0
	# 创建一个workbook 设置编码
	workbook = xlwt.Workbook(encoding='utf-8')
	# 创建一个worksheet
	worksheet = workbook.add_sheet('My Worksheet')
	
	# 写入excel
	# 参数对应 行, 列, 值
	worksheet.write(0, 0, label='朗读大纲')
	worksheet.write(1, 0, label='页号')
	worksheet.write(1, 1, label='行号')
	worksheet.write(1, 2, label='文本内容')
	for slide in prs.slides:
		picture_num = picture_num + 1
		line_num = 0
		# 获取形状shape
		for shape in slide.shapes:
			
			if shape.has_text_frame:  # 判断是否有文字
				
				text_frame = shape.text_frame.text  # 获取文字框
				if len(text_frame) != 0:
					count_num = count_num + 1
					line_num = line_num + 1
					
					worksheet.write(count_num + 1, 0, label=str(picture_num))
					worksheet.write(count_num + 1, 1, label=str(line_num))
					worksheet.write(count_num + 1, 2, label=text_frame)
	workbook.save('朗读1.xls')


def check_directory():
	'''
	检查文件目录结构
	:return:
	'''
	if not os.path.isfile('朗读.pptx'):
		print('请另存待处理文件为‘朗读.pptx’然后重新运行此程序。')
		input('按任意键退出此程序')
	if os.path.isdir('./images'):
		shutil.rmtree('images')
	if os.path.isdir('./audios'):
		shutil.rmtree('audios')
	if os.path.isdir('./temp'):
		shutil.rmtree('temp')
	os.mkdir('images')
	os.mkdir('audios')
	os.mkdir('temp')


def xls2mp3():
	'''
	朗读xls文件，转为MP3
	:return:
	'''
	book = xlrd.open_workbook('朗读1.xls')
	
	sheet1 = book.sheets()[0]
	
	nrows = sheet1.nrows
	temp_text = ''
	print('表格总行数', nrows)
	temp_row = 1
	for i in range(2, nrows):
		reading(sheet1.cell(i, 2).value,
		        sheet1.cell(i, 0).value + '-' + sheet1.cell(i, 1).value + "-" + sheet1.cell(i, 2).value[:5])
		if sheet1.cell(i, 0).value == temp_row:
			temp_text = temp_text + "。。。。" + sheet1.cell(i, 2).value
		else:
			if temp_row == 1:
				temp_text = temp_text + "。。。。" + sheet1.cell(i, 2).value
			else:
				# print('读', temp_text)
				reading(temp_text + "。。。。", temp_row)
			temp_row = sheet1.cell(i, 0).value
			temp_text = sheet1.cell(i, 2).value
	
	# print('读', temp_text)
	reading(temp_text + "。。。。", temp_row)


# reading(temp_text + "。。。。", int(sheet1.cell(i, 0).value))
# else:
# 	temp_text = temp_text + "。。。。" + sheet1.cell(i, 2).value


# reading(temp_text + "。。。。", int(sheet1.cell(i, 0).value))


# 	if sheet1.cell(i, 0).value == temp_row:
# 		temp_text = temp_text +"。。。。"+ sheet1.cell(i, 2).value
# 	else:
# 		reading(temp_text+"。。。。", int(sheet1.cell(i, 0).value) - 1)
# 		temp_row = sheet1.cell(i, 0).value
# 		temp_text = sheet1.cell(i, 2).value
# reading(temp_text, sheet1.cell(i, 0).value)


# reading(sheet1.cell(i, 2).value, sheet1.cell(i, 0).value + '-' + sheet1.cell(i, 1).value)


def audio_duration(file):
	# print(file)
	
	audioclip = AudioFileClip(file)
	# print(audioclip)
	
	s = audioclip.duration
	# print(s)
	s = math.ceil(s)
	if s < 5:
		return 5
	return s


def merge_video_audio():
	'''
	合并音视频
	:return:
	'''
	all_clip = []
	all_file = os.listdir('images')
	# print(len(all_file))
	for i in range(len(all_file)):
		auido = 'audios/%s.mp3' % (i + 1)
		image = 'images/幻灯片%s.JPG' % (i + 1)
		print(auido, image)
		duration = audio_duration(auido)
		# print(duration)
		
		v_clip = ImageSequenceClip([image, ], durations=[duration])
		a_clip = AudioFileClip(auido)
		resout = v_clip.set_audio(a_clip)
		all_clip.append(resout)
		resout.write_videofile('./temp/%s.mp4' % i, fps=24)
	finalclip = concatenate_videoclips(all_clip)
	finalclip.write_videofile("my_concatenate.mp4", fps=24)


def stop():
	input('现在是修正读音暂停，请校验朗读.xls后按任意键继续')


def merge_all():
	all_file = os.listdir('temp')
	# print(len(all_file))
	for i in range(1, len(all_file) + 1):
		print(i)


def main():
	check_directory()
	ppt2jpg()
	ppt2xls()
	stop()
	xls2mp3()
	merge_video_audio()
	# check_directory()


# merge_all()


if __name__ == '__main__':
	main()
